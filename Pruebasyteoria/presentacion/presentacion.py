from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Crear la presentación
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]
section_slide_layout = prs.slide_layouts[2]

# Función para añadir una diapositiva con título y bullets
def add_bullet_slide(title, bullet_points):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = title
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0

# Portada
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Marco Jurídico y Oportunidades de Negocio en el Software Libre"
slide.placeholders[1].text = "Basado en el artículo de Jordi Mas"

# Secciones de la presentación
add_bullet_slide("1. Industria del Software y el Software Libre", [
    "El software libre se ha convertido en una herramienta estratégica.",
    "Ejemplos: Linux, Apache, Firefox.",
    "Transforma el modelo de negocio tradicional.",
    "El 70% de los ingresos provienen de servicios asociados."
])

add_bullet_slide("2. Marco Legal del Software Libre", [
    "2.1 Derechos de autor: protege el código fuente, no las ideas.",
    "2.2 Licencias: Copyleft (GPL), Permisivas (MIT, BSD), Dual.",
    "2.3 Patentes: controvertidas, especialmente en Europa.",
    "2.4 Marcas: el software es libre, la marca puede no serlo."
])

add_bullet_slide("3. Modelos de Negocio", [
    "Venta de software con valor añadido (ej. Red Hat).",
    "Servicios profesionales: soporte, formación, personalización.",
    "Integración en hardware (ej. Android en smartphones)."
])

add_bullet_slide("4. Casos Reales", [
    "Fracasos empresariales: Abiword, Eazel.",
    "Casos de éxito: Red Hat, IBM, MySQL, Ándago.",
    "Las comunidades aseguran la continuidad del software."
])

add_bullet_slide("5. Conclusiones", [
    "El software libre es viable y estratégico.",
    "Ofrece independencia, innovación y reducción de costos.",
    "Fortalece el ecosistema tecnológico local y global."
])

# Guardar el archivo
pptx_path = "/mnt/data/Presentacion_SW_Libre.pptx"
prs.save(pptx_path)
pptx_path
